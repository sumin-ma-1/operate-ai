import base64
import io
import subprocess
import sys
import tempfile
from pathlib import Path


def extract_document_text(filename: str, content_base64: str) -> str:
    extension = Path(filename).suffix.lower()
    data = base64.b64decode(content_base64)

    extractors = {
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".pdf": _extract_pdf,
        ".hwp": _extract_hwp,
    }

    extractor = extractors.get(extension)
    if extractor is None:
        supported = ", ".join(sorted(extractors))
        raise ValueError(
            f"Unsupported document type: {filename}. Supported: {supported}"
        )

    text = extractor(data).strip()
    if not text:
        raise ValueError(f"No readable text found in {filename}.")

    return text


def _extract_docx(data: bytes) -> str:
    from docx import Document

    document = Document(io.BytesIO(data))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n".join(lines)


def _extract_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []

    for sheet in workbook.worksheets:
        lines.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(cell.strip() for cell in cells):
                lines.append("\t".join(cells))

    workbook.close()
    return "\n".join(lines)


def _extract_pptx(data: bytes) -> str:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    lines: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"## Slide {index}")
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())

    return "\n".join(lines)


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(page.strip() for page in pages if page.strip())


def _extract_hwp(data: bytes) -> str:
    with tempfile.NamedTemporaryFile(suffix=".hwp", delete=False) as temp_file:
        temp_file.write(data)
        temp_path = temp_file.name

    try:
        result = subprocess.run(
            [sys.executable, "-m", "hwp5.hwp5txt", temp_path],
            capture_output=True,
            text=True,
            timeout=30,
            check=False,
        )
    finally:
        Path(temp_path).unlink(missing_ok=True)

    if result.returncode != 0:
        detail = (result.stderr or result.stdout or "Unknown HWP parser error").strip()
        raise ValueError(
            "Could not read HWP file. Try exporting to .docx or .pdf from Hancom Office. "
            f"({detail})"
        )

    return result.stdout
